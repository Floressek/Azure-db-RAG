import os
import io
import csv
import json
import xml.etree.ElementTree as ET
from pptx import Presentation
from docx import Document
from PyPDF2 import PdfReader
import openpyxl
import pandas as pd


def parse_pptx(file_path):
    presentation = Presentation(file_path)
    slides = []
    for i, slide in enumerate(presentation.slides, 1):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                slide_content.append(shape.text)
        slides.append({f"Slide {i}": slide_content})
    return {"slides": slides}


def parse_doc(file_path):
    document = Document(file_path)
    paragraphs = [paragraph.text for paragraph in document.paragraphs]
    return {"paragraphs": paragraphs}


def parse_pdf(file_path):
    reader = PdfReader(file_path)
    pages = []
    for i, page in enumerate(reader.pages, 1):
        pages.append({f"Page {i}": page.extract_text()})
    return {"pages": pages}


def parse_xml(file_path):
    tree = ET.parse(file_path)
    root = tree.getroot()

    def element_to_dict(element):
        result = {}
        for child in element:
            child_data = element_to_dict(child)
            if child.tag in result:
                if type(result[child.tag]) is list:
                    result[child.tag].append(child_data)
                else:
                    result[child.tag] = [result[child.tag], child_data]
            else:
                result[child.tag] = child_data
        if element.text and element.text.strip():
            if len(result) > 0:
                result['#text'] = element.text.strip()
            else:
                result = element.text.strip()
        return result

    return element_to_dict(root)


def parse_excel(file_path):
    workbook = openpyxl.load_workbook(file_path)
    sheets = {}
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        data = []
        for row in sheet.iter_rows(values_only=True):
            data.append(list(row))
        sheets[sheet_name] = data
    return {"sheets": sheets}


def parse_csv(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        reader = csv.reader(file)
        data = list(reader)
    return {"data": data}


def parse_file(file_path):
    _, file_extension = os.path.splitext(file_path)
    file_extension = file_extension.lower()

    if file_extension == '.pptx':
        return parse_pptx(file_path)
    elif file_extension in ['.doc', '.docx']:
        return parse_doc(file_path)
    elif file_extension == '.pdf':
        return parse_pdf(file_path)
    elif file_extension == '.xml':
        return parse_xml(file_path)
    elif file_extension in ['.xlsx', '.xls']:
        return parse_excel(file_path)
    elif file_extension == '.csv':
        return parse_csv(file_path)
    else:
        raise ValueError(f"Unsupported file format: {file_extension}")


# Example usage
file_path = r'C:\Users\szyme\PycharmProjects\pythonProject2\pdf_parser\test1.pdf'
parsed_content = parse_file(file_path)
json_output = json.dumps(parsed_content, ensure_ascii=False, indent=2)
print(json_output)

# Optionally, save to a file
with open('output_v1.json', 'w', encoding='utf-8') as f:
    json.dump(parsed_content, f, ensure_ascii=False, indent=2)
